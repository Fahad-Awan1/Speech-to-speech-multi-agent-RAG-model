import os
import io
import time
import tempfile
from pathlib import Path
import hashlib

import streamlit as st
from audio_recorder_streamlit import audio_recorder
from pydub import AudioSegment
import requests

# ----- Core Logic Imports -----
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from groq import Groq
from dotenv import load_dotenv
import docx2txt
from pypdf import PdfReader
from pptx import Presentation
from rank_bm25 import BM25Okapi

# ----- TTS Imports -----
USE_COQUI = True
try:
    from TTS.api import TTS
except ImportError:
    st.warning(
        "Coqui TTS not found, falling back to gTTS. For better quality, run: pip install TTS"
    )
    USE_COQUI = False
from gtts import gTTS

# ---------------------------
# Load Environment & Page Config
# ---------------------------
load_dotenv()

st.set_page_config(
    page_title="AI Assistant",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --- Check for API Keys on Startup ---
if not os.getenv("ASSEMBLYAI_API_KEY") or not os.getenv("GROQ_API_KEY"):
    st.error(
        "🚨 API Key Error: Please ensure both ASSEMBLYAI_API_KEY and GROQ_API_KEY are set in your .env file."
    )
    st.stop()

# ---------------------------
# Caching for Performance
# ---------------------------


@st.cache_resource
def load_sentence_transformer_model(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
):
    return SentenceTransformer(model_name)


@st.cache_resource
def load_tts_model():
    if USE_COQUI:
        try:
            return TTS(model_name="tts_models/en/ljspeech/tacotron2-DDC")
        except Exception as e:
            st.error(f"Failed to load Coqui TTS model: {e}")
            return None
    return None


# ---------------------------
# Enhanced Backend Functionality (Unchanged)
# ---------------------------


def load_text_from_filelike(filename: str, data: bytes) -> str:
    # This function remains the same as the previous version
    suffix = Path(filename).suffix.lower()
    text = ""
    try:
        if suffix == ".pdf":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(data)
                tmp.flush()
                reader = PdfReader(tmp.name)
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif suffix == ".docx":
            with io.BytesIO(data) as docx_file:
                text = docx2txt.process(docx_file) or ""
        elif suffix == ".pptx":
            with io.BytesIO(data) as pptx_file:
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        elif suffix in (".txt", ".md"):
            text = data.decode("utf-8", errors="ignore")
    except Exception as e:
        st.error(f"Error processing {filename}: {e}")
    return text


def chunk_text(text: str, size=1000, overlap=150):
    # This function remains the same
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += size - overlap
    return chunks


class HybridRetriever:
    # This class remains the same
    def __init__(self, model):
        self.embedder = model
        self.texts = []
        self.sources = []
        self.faiss_index = None
        self.bm25_index = None
        self.dim = self.embedder.get_sentence_embedding_dimension()

    def build(self, docs: list[tuple[str, str]]):
        self.texts = [t for t, _ in docs]
        self.sources = [s for _, s in docs]
        embeddings = self.embedder.encode(
            self.texts, normalize_embeddings=True, convert_to_numpy=True
        )
        self.faiss_index = faiss.IndexFlatIP(self.dim)
        self.faiss_index.add(embeddings)
        tokenized_corpus = [doc.lower().split() for doc in self.texts]
        self.bm25_index = BM25Okapi(tokenized_corpus)

    def search(self, query: str, k=5):
        if not self.texts:
            return []
        query_embedding = self.embedder.encode(
            [query], normalize_embeddings=True, convert_to_numpy=True
        )
        faiss_scores, faiss_indices = self.faiss_index.search(query_embedding, k * 5)
        tokenized_query = query.lower().split()
        bm25_scores = self.bm25_index.get_scores(tokenized_query)
        bm25_top_indices = np.argsort(bm25_scores)[::-1][: k * 5]
        fused_scores = {}
        rrf_k = 60
        for rank, idx in enumerate(faiss_indices[0]):
            fused_scores[idx] = fused_scores.get(idx, 0) + 1.0 / (rank + rrf_k)
        for rank, idx in enumerate(bm25_top_indices):
            fused_scores[idx] = fused_scores.get(idx, 0) + 1.0 / (rank + rrf_k)
        sorted_indices = sorted(
            fused_scores.keys(), key=lambda idx: fused_scores[idx], reverse=True
        )
        return [
            {"text": self.texts[idx], "source": self.sources[idx]}
            for idx in sorted_indices[:k]
        ]


class Services:
    # This class remains the same
    def __init__(self, model_name):
        self.groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        self.model_name = model_name
        self.tts_model = load_tts_model()

    def transcribe_audio(self, wav_bytes: bytes) -> str:
        api_key = os.getenv("ASSEMBLYAI_API_KEY")
        headers = {"authorization": api_key}
        base_url = "https://api.assemblyai.com/v2"
        try:
            upload_response = requests.post(
                f"{base_url}/upload", headers=headers, data=wav_bytes
            )
            upload_response.raise_for_status()
            audio_url = upload_response.json()["upload_url"]
            transcript_response = requests.post(
                f"{base_url}/transcript", json={"audio_url": audio_url}, headers=headers
            )
            transcript_response.raise_for_status()
            transcript_id = transcript_response.json()["id"]
            polling_endpoint = f"{base_url}/transcript/{transcript_id}"
            while True:
                result = requests.get(polling_endpoint, headers=headers).json()
                if result["status"] == "completed":
                    return result["text"] or ""
                elif result["status"] == "error":
                    raise RuntimeError(f"Transcription failed: {result['error']}")
                time.sleep(2)
        except requests.exceptions.RequestException as e:
            st.error(f"AssemblyAI API Error: {e}")
            return ""

    def generate_streamed_response(
        self, query: str, contexts: list[dict], chat_history: list
    ):
        context_block = "\n\n".join(
            [f"Source: {c['source']}\nContent: {c['text']}" for c in contexts]
        )
        history_block = "\n".join(
            [f"{msg['role']}: {msg['content']}" for msg in chat_history]
        )
        prompt = f"You are a helpful assistant. Answer the user's question based *only* on the provided context. If the answer is not in the context, state that you cannot find the information. \n\nPrevious Conversation:\n{history_block}\n\nProvided Context:\n{context_block}\n\nUser Question: {query}\n\nAnswer:"
        try:
            stream = self.groq_client.chat.completions.create(
                model=self.model_name,
                messages=[{"role": "user", "content": prompt}],
                stream=True,
            )
            for chunk in stream:
                yield chunk.choices[0].delta.content or ""
        except Exception as e:
            st.error(f"Groq API Error: {e}")
            yield ""

    def synthesize_audio(self, text: str) -> bytes | None:
        try:
            if USE_COQUI and self.tts_model:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
                    self.tts_model.tts_to_file(text=text, file_path=tmp.name)
                    return Path(tmp.name).read_bytes()
            else:
                tts = gTTS(text=text, lang="en")
                with io.BytesIO() as fp:
                    tts.write_to_fp(fp)
                    fp.seek(0)
                    mp3_audio = AudioSegment.from_file(fp, format="mp3")
                    wav_io = io.BytesIO()
                    mp3_audio.export(wav_io, format="wav")
                    return wav_io.getvalue()
        except Exception as e:
            st.error(f"Text-to-Speech Error: {e}")
            return None


# ---------------------------
# Streamlit UI
# ---------------------------

# --- NEW CSS FOR THE DARK THEME ---
st.markdown(
    """
<style>
    /* General App Styling */
    .stApp { background-color: #0F1115; color: #E5E7EB; }
    #MainMenu, footer, header { visibility: hidden; }
    .block-container { padding: 1rem 2rem 2rem 2rem !important; }

    /* Top Gradient Bar */
    .top-gradient-bar {
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 4px;
        background: linear-gradient(90deg, #ff4e50, #f9d423);
        z-index: 999999;
    }

    /* Sidebar Styling */
    [data-testid="stSidebar"] {
        background-color: #1F2125;
        padding: 1rem;
    }
    [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2 {
        color: #FFFFFF;
        font-weight: 300;
    }
    [data-testid="stSidebar"] p {
        color: #A0A0A0;
    }
    .stButton > button {
        background-color: #30333A;
        color: #FFFFFF;
        border: 1px solid #4A4D55;
        border-radius: 8px;
        padding: 10px 14px;
    }
    .stButton > button:hover {
        background-color: #3C3F44;
        border-color: #5A5F69;
    }
    [data-testid="stSidebar"] .stButton > button {
        background-color: #30333A;
        color: #FFFFFF;
        border: 1px solid #4A4D55;
        border-radius: 8px;
        padding: 10px 0;
    }

    /* Custom File Uploader */
    [data-testid="stFileUploader"] {
        background-color: #1E1F22;
        border: 1px solid #3C3F44;
        border-radius: 8px;
        padding: 1.5rem;
    }
    [data-testid="stFileUploader"] > label {
        color: #FFFFFF;
        font-weight: 300;
    }
    [data-testid="stFileUploader"] small {
        color: #A0A0A0;
    }
    [data-testid="stFileUploader"] button {
        background-color: #3C3F44;
        color: #FFFFFF;
        border: 1px solid #555;
    }
    
    /* Main Content Area */
    .main-header {
        color: #E0E0E0;
        font-weight: 200;
        font-size: 3rem;
        text-align: center;
        margin-bottom: 2rem;
    }

    /* Input Fields */
    [data-testid="stTextInput"] > div > div > input, .st-emotion-cache-1ypb00x {
        background-color: #25262A !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 15px !important;
        box-shadow: none !important;
    }

    /* Chat Bubbles (Simple styling for when chat history appears) */
    .user-bubble { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 1rem; border-radius: 10px 10px 0 10px; }
    .assistant-bubble { background-color: #2B2D31; color: #E5E7EB; padding: 1rem; border-radius: 10px 10px 10px 0; border: 1px solid #3C3F44; }

</style>
""",
    unsafe_allow_html=True,
)

# Add the gradient bar to the top of the page
st.markdown('<div class="top-gradient-bar"></div>', unsafe_allow_html=True)

# --- Session State ---
if "retriever" not in st.session_state:
    st.session_state.retriever = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "processing" not in st.session_state:
    st.session_state.processing = False

# ---------------------------
# Sidebar UI
# ---------------------------
with st.sidebar:
    st.title("🚀 AI Assistant")
    st.markdown("---")
    st.header("1. Upload Documents")
    st.caption("PDF, DOCX, PPTX, TXT")
    files = st.file_uploader(
        "Drag and drop files here",
        type=["pdf", "docx", "pptx", "txt", "md"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )
    if st.button("Build Knowledge Base", use_container_width=True):
        if not files:
            st.warning("Please upload at least one document.")
        else:
            with st.spinner("Processing documents..."):
                docs = [
                    (ch, f.name)
                    for f in files
                    if (text := load_text_from_filelike(f.name, f.read())).strip()
                    for ch in chunk_text(text)
                ]
                if docs:
                    retriever = HybridRetriever(load_sentence_transformer_model())
                    retriever.build(docs)
                    st.session_state.retriever = retriever
                    st.success(f"✅ Indexed {len(retriever.texts)} text chunks!")
                else:
                    st.error("Could not extract any text.")

    st.markdown("---")
    st.header("2. Configure Settings")
    top_k = st.slider("Context Chunks", 1, 15, 5)
    llm_model = st.selectbox(
        "Language Model", ["llama3-70b-8192", "llama3-8b-8192", "mixtral-8x7b-32768"]
    )
    st.divider()
    if st.button("Clear Chat History", use_container_width=True):
        st.session_state.chat_history = []
        st.rerun()

# ---------------------------
# Main Chat UI
# ---------------------------
st.markdown(
    '<h1 class="main-header">Enhanced AI Document Assistant</h1>',
    unsafe_allow_html=True,
)

# --- Chat History Display (appears once there's history) ---
if st.session_state.chat_history:
    for message in st.session_state.chat_history:
        if message["role"] == "user":
            st.markdown(
                f'<div class="user-bubble">{message["content"]}</div>',
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                f'<div class="assistant-bubble">{message["content"]}</div>',
                unsafe_allow_html=True,
            )
            if message.get("audio"):
                st.audio(message["audio"], format="audio/wav")
            if message.get("context"):
                with st.expander("Show Sources"):
                    for i, ctx in enumerate(message["context"]):
                        st.info(f"Source {i + 1}: {ctx['source']}\n\n> {ctx['text']}")
    st.markdown("---")


# --- User Input ---
text_query = st.text_input(
    "Type your question...",
    key="text_query",
    disabled=st.session_state.processing,
    label_visibility="collapsed",
)

audio_bytes = None
# The custom class st-emotion-cache-1ypb00x is what audio_recorder creates for its button
# We target it to apply our custom dark theme styling.
if not st.session_state.processing:
    audio_bytes = audio_recorder(
        text="Click mic to ask...", icon_size="2x", key="voice_input"
    )
else:
    st.button("🎤 Processing...", disabled=True, use_container_width=True)

query = text_query if text_query else ("audio" if audio_bytes else None)

# --- Processing Logic  ---
if query and not st.session_state.processing:
    if not st.session_state.retriever:
        st.error("Please build a knowledge base first.")
    else:
        st.session_state.processing = True
        query_text = ""
        try:
            if query == "audio":
                with st.spinner("🎤 Transcribing..."):
                    services = Services(llm_model)
                    query_text = services.transcribe_audio(audio_bytes)
            else:
                query_text = query
            if query_text.strip():
                st.session_state.chat_history.append(
                    {"role": "user", "content": query_text}
                )
                st.rerun()
        except Exception as e:
            st.error(f"Input processing error: {e}")
            st.session_state.processing = False

if (
    st.session_state.chat_history
    and st.session_state.chat_history[-1]["role"] == "user"
):
    st.session_state.processing = True
    last_user_message = st.session_state.chat_history[-1]["content"]
    services = Services(llm_model)
    retriever = st.session_state.retriever

    with st.spinner("🔍 Searching documents..."):
        context = retriever.search(last_user_message, k=top_k)

    if not context:
        st.warning("Could not find relevant information for your query.")
        st.session_state.processing = False
    else:
        with st.chat_message("assistant", avatar="🚀"):
            # 2. Call write_stream INSIDE the container. It will render the stream here.
            full_response = st.write_stream(services.generate_streamed_response(...))

        with st.spinner("🎧 Synthesizing audio..."):
            audio_response = services.synthesize_audio(full_response)

        st.session_state.chat_history.append(
            {
                "role": "assistant",
                "content": full_response,
                "audio": audio_response,
                "context": context,
            }
        )
        st.session_state.processing = False
        st.rerun()
